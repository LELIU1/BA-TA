#!/usr/bin/env python3
"""生成字节跳动BA画像PPT v3.3 — 严格遵循设计规范"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── 全局常量 ────────────────────────────────────────────
BLUE    = RGBColor(0x25, 0x63, 0xEB)
DGRAY   = RGBColor(0x1E, 0x29, 0x3B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF1, 0xF5, 0xF9)
LBLUE   = RGBColor(0xEF, 0xF6, 0xFF)
BLACK   = RGBColor(0x00, 0x00, 0x00)
FONT    = 'PingFang SC'
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.3)
TITLE_H  = Inches(0.7)
CONTENT_TOP = Inches(0.8)
CONTENT_BOT = Inches(7.15)
GAP     = Inches(0.15)

# ── 辅助函数 ────────────────────────────────────────────

def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, text):
    """添加浅灰标题栏"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, TITLE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = LGRAY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = MARGIN
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DGRAY
    p.font.name = FONT
    return bar

def add_textbox(slide, left, top, width, height, text="", font_size=Pt(11),
                bold=False, color=DGRAY, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                font_name=FONT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox

def add_card(slide, left, top, width, height, fill_color=WHITE, border_color=LGRAY):
    """添加白底卡片"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(0.5)
    return card

def add_blue_left_border(slide, left, top, height, width=Inches(0.06)):
    """添加左侧深蓝色竖线"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    return line

def set_cell_text(cell, text, font_size=Pt(11), bold=False, color=DGRAY, alignment=PP_ALIGN.LEFT, font_name=FONT):
    """设置表格单元格文字"""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 减小内边距
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)

def style_table_header(table):
    """表头深蓝底白字"""
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.size = Pt(11)

def style_table_rows(table, start_row=1):
    """交替白/浅灰背景"""
    for i in range(start_row, len(table.rows)):
        row = table.rows[i]
        for cell in row.cells:
            cell.fill.solid()
            if (i - start_row) % 2 == 0:
                cell.fill.fore_color.rgb = WHITE
            else:
                cell.fill.fore_color.rgb = LGRAY

def add_arrow_shape(slide, left, top, width, height):
    """添加向下箭头"""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLUE
    arrow.line.fill.background()
    return arrow

def add_flow_box(slide, left, top, width, height, text, font_size=Pt(10), fill=BLUE, text_color=WHITE):
    """添加深蓝色流程框"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    # 调整圆角
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = text_color
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return box

def add_kpi_card(slide, left, top, width, height, label, value):
    """添加KPI数字卡片"""
    card = add_card(slide, left, top, width, height, fill_color=LBLUE, border_color=LGRAY)
    # 数值
    add_textbox(slide, left + Inches(0.05), top + Inches(0.05), width - Inches(0.1), Inches(0.35),
                value, font_size=Pt(14), bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    # 标签
    add_textbox(slide, left + Inches(0.05), top + Inches(0.38), width - Inches(0.1), Inches(0.3),
                label, font_size=Pt(9), bold=False, color=DGRAY, alignment=PP_ALIGN.CENTER)

def set_slide_bg_all(prs):
    for slide in prs.slides:
        set_slide_bg(slide)

# ═══════════════════════════════════════════════════════════
#  创建 PPT
# ═══════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# 空白版式
blank_layout = prs.slide_layouts[6]  # blank

# ═══════════════════════════════════════════════════════════
#  第1页：客户画像 + 商业模式
# ═══════════════════════════════════════════════════════════

slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1)
add_title_bar(slide1, "字节跳动 BA画像 — 客户基本画像与商业模式")

# ── 上半部分：客户基本信息表 + KPI卡片 ──
top_y = CONTENT_TOP

# 左侧：客户基本信息表 (3行)
tbl1 = slide1.shapes.add_table(4, 2, MARGIN, top_y, Inches(5.8), Inches(1.6)).table
tbl1.columns[0].width = Inches(1.6)
tbl1.columns[1].width = Inches(4.2)

info_data = [
    ["维度", "内容"],
    ["客户名称", "字节跳动（ByteDance）"],
    ["所属行业", "互联网/科技 — 社交媒体、内容平台、云计算、AI"],
    ["营收规模", "2025年 ~1860亿美元；净利润~500亿美元（AI投入压缩利润）"]
]
for r, row_data in enumerate(info_data):
    for c, val in enumerate(row_data):
        set_cell_text(tbl1.cell(r, c), val, bold=(r == 0))
style_table_header(tbl1)
style_table_rows(tbl1)

# 右侧：5张KPI卡片并排
kpi_x = Inches(6.3)
kpi_w = Inches(1.35)
kpi_h = Inches(0.7)

kpis = [
    ("豆包MAU", "3.45亿"),
    ("日均Token", "120万亿"),
    ("Capex上限", "$700亿"),
    ("MaaS份额", "49.5%"),
    ("Trae开发者", "600万"),
]

for i, (label, value) in enumerate(kpis):
    x = kpi_x + i * (kpi_w + GAP)
    add_kpi_card(slide1, x, top_y, kpi_w, kpi_h, label, value)

# 发展阶段文字
add_textbox(slide1, kpi_x, top_y + kpi_h + Inches(0.1), Inches(6.7), Inches(0.7),
            "发展阶段：成熟期 + AI转型深水区\n"
            "核心产品三线增长：内容平台(AI短剧/漫剧)、AIGC应用(豆包付费)、火山引擎(Agent Plan)",
            font_size=Pt(10), color=DGRAY)

# ── 下半部分 ──
lower_y = top_y + Inches(1.85)
lower_h = CONTENT_BOT - lower_y

# 左侧：核心业务矩阵表 (6类)
matrix_title_y = lower_y
add_textbox(slide1, MARGIN, matrix_title_y, Inches(5.8), Inches(0.3),
            "核心业务矩阵", font_size=Pt(12), bold=True, color=BLUE)

tbl2 = slide1.shapes.add_table(8, 4, MARGIN, matrix_title_y + Inches(0.3),
                                Inches(6.3), Inches(4.2)).table
tbl2.columns[0].width = Inches(1.2)
tbl2.columns[1].width = Inches(1.8)
tbl2.columns[2].width = Inches(2.0)
tbl2.columns[3].width = Inches(1.3)

biz_data = [
    ["类别", "产品", "规模/地位", "AI算力关联"],
    ["内容与社区", "抖音", "DAU 7亿+", "推荐LLM化↑↑"],
    ["", "TikTok", "全球DAU 10亿+", "海外AIDC"],
    ["", "番茄/红果短剧", "MAU 2.45/3.04亿", "视频推理↑↑↑"],
    ["创作工具", "剪映/即梦AI", "全球TOP1 / MAU 7994万", "Seedance推理"],
    ["AI应用", "豆包App", "MAU 3.45亿", "最大推理消费者"],
    ["", "Trae/Coze", "MAU 600万+", "Agent 7×24"],
    ["企业服务", "火山引擎/飞书", "MaaS 49.5%第一", "全栈AI输出"],
]
for r, row_data in enumerate(biz_data):
    for c, val in enumerate(row_data):
        set_cell_text(tbl2.cell(r, c), val, font_size=Pt(9.5), bold=(r == 0))
style_table_header(tbl2)
style_table_rows(tbl2)

# 右侧：地域结构表 + 收入结构表
right_x = Inches(6.7)
right_w = Inches(6.3)

add_textbox(slide1, right_x, lower_y, right_w, Inches(0.3),
            "地域结构", font_size=Pt(12), bold=True, color=BLUE)

tbl3 = slide1.shapes.add_table(4, 4, right_x, lower_y + Inches(0.3),
                                right_w, Inches(1.5)).table
tbl3.columns[0].width = Inches(1.2)
tbl3.columns[1].width = Inches(1.2)
tbl3.columns[2].width = Inches(1.2)
tbl3.columns[3].width = Inches(2.7)

geo_data = [
    ["区域", "营收占比", "增速", "数据源"],
    ["国内", "~70%", "~20% YoY", "36氪独家(2026.04)"],
    ["海外", "~30%", "~50% YoY", "TikTok Shop GMV增速~70%"],
    ["合计", "100%", "整体高增长", "驱动：海外营收占比突破30%"],
]
for r, row_data in enumerate(geo_data):
    for c, val in enumerate(row_data):
        set_cell_text(tbl3.cell(r, c), val, font_size=Pt(9.5), bold=(r == 0))
style_table_header(tbl3)
style_table_rows(tbl3)

# 收入结构表
rev_top = lower_y + Inches(2.0)
add_textbox(slide1, right_x, rev_top, right_w, Inches(0.3),
            "收入结构", font_size=Pt(12), bold=True, color=BLUE)

tbl4 = slide1.shapes.add_table(5, 2, right_x, rev_top + Inches(0.3),
                                right_w, Inches(1.7)).table
tbl4.columns[0].width = Inches(2.5)
tbl4.columns[1].width = Inches(3.8)

rev_data = [
    ["来源", "占比与说明"],
    ["广告", "~55-60% — 抖音/TikTok/头条信息流"],
    ["电商", "~15-20% — 抖音电商+TikTok Shop"],
    ["云+AI", "~5-10% — 火山引擎MaaS+豆包付费"],
    ["游戏/教育/其他", "~5-10%"],
]
for r, row_data in enumerate(rev_data):
    for c, val in enumerate(row_data):
        set_cell_text(tbl4.cell(r, c), val, font_size=Pt(9.5), bold=(r == 0))
style_table_header(tbl4)
style_table_rows(tbl4)


# ═══════════════════════════════════════════════════════════
#  第2页：核心业务流程 — 双轨制
# ═══════════════════════════════════════════════════════════

slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2)
add_title_bar(slide2, "核心业务流程 — 双轨制：搜索推荐AI闭环 vs AIGC-Agent应用闭环")

top_y2 = CONTENT_TOP
left_w = Inches(6.3)
right_w2 = Inches(6.3)
right_x2 = Inches(6.8)
flow_box_w = Inches(5.5)

# ── 左侧：搜索推荐AI闭环 ──
add_textbox(slide2, MARGIN, top_y2, left_w, Inches(0.35),
            "左轨：搜索推荐AI闭环（传统核心）", font_size=Pt(12), bold=True, color=BLUE)

ly = top_y2 + Inches(0.4)
box_h = Inches(0.65)
arrow_h = Inches(0.25)
gap_v = Inches(0.08)

# 流程图 - 用深蓝色
steps_left = [
    "用户行为数据\n抖音/TikTok/头条海量行为",
    "LLM级生成式推荐\nRankMixer→OneTrans→广告电商转化",
    "AI基础设施投入\n现金流支撑Capex",
    "更强模型→更好体验\n→更多数据→封闭飞轮",
]

for i, txt in enumerate(steps_left):
    bh = box_h if i != 1 else Inches(0.85)  # 第二个框更高
    y = ly + i * (box_h + arrow_h + gap_v)
    add_flow_box(slide2, MARGIN + Inches(0.4), y, flow_box_w, bh, txt, font_size=Pt(9.5))
    if i < len(steps_left) - 1:
        ay = y + bh
        add_arrow_shape(slide2, MARGIN + Inches(0.4) + flow_box_w/2 - Inches(0.1),
                        ay, Inches(0.2), arrow_h)

# 芯片需求标注
chip_y = ly + 4 * (box_h + arrow_h + gap_v) - Inches(0.05)
add_textbox(slide2, MARGIN, chip_y, left_w, Inches(0.4),
            "⚙ 芯片需求：训练GPU（昇腾910C + AMD）为主",
            font_size=Pt(10), bold=False, color=DGRAY)

# ── 右侧：AIGC-Agent应用闭环 ──
add_textbox(slide2, right_x2, top_y2, right_w2, Inches(0.35),
            "右轨：AIGC-Agent应用闭环（新AI增长极）⭐ 主线", font_size=Pt(12), bold=True, color=BLUE)

ry = top_y2 + Inches(0.4)
right_box_w = Inches(1.4)
right_box_h = Inches(1.1)
right_gap = Inches(0.12)

# 上层4个应用框
app_boxes = [
    "豆包App\nMAU 3.45亿",
    "即梦+剪映\nSeedance 2.0",
    "Trae+Coze\n600万开发者",
    "火山引擎MaaS\n49.5%份额",
]
n_apps = len(app_boxes)
total_w = n_apps * right_box_w + (n_apps - 1) * right_gap
start_x = right_x2 + (right_w2 - total_w) / 2

for i, txt in enumerate(app_boxes):
    x = start_x + i * (right_box_w + right_gap)
    add_flow_box(slide2, x, ry, right_box_w, right_box_h, txt, font_size=Pt(8.5))

# 中间层：算法策略与模型矩阵
mid_y = ry + right_box_h + arrow_h + Inches(0.05)
add_flow_box(slide2, right_x2 + Inches(0.4), mid_y, Inches(5.5), Inches(0.55),
             "算法策略与模型矩阵：Seed-2.0-pro(200B MoE) + Seed-Thinking-v1.5 + lite全模态",
             font_size=Pt(9.5))

# 下层：AI芯片与算力供应
bottom_y = mid_y + Inches(0.55) + arrow_h + Inches(0.05)
add_flow_box(slide2, right_x2 + Inches(0.4), bottom_y, Inches(5.5), Inches(0.55),
             "AI芯片与算力供应：昇腾950PR + 寒武纪590 + 高通ASIC + 自研Ada",
             font_size=Pt(9.5))

# 底部箭头
final_arrow_y = bottom_y + Inches(0.55)
add_arrow_shape(slide2, right_x2 + right_w2/2 - Inches(0.1), final_arrow_y, Inches(0.2), arrow_h)

# 底部标注
final_y = final_arrow_y + arrow_h + Inches(0.05)
add_textbox(slide2, right_x2, final_y, right_w2, Inches(0.45),
            "日均120万亿Token × Agent 7×24 × 百万Token上下文 → 推理芯片扩容",
            font_size=Pt(10), bold=True, color=BLUE)

# ── 页面底部窄条：核心业务痛点 ──
pain_y = Inches(6.7)
pain_bar = slide2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, MARGIN, pain_y, Inches(12.7), Inches(0.65))
pain_bar.fill.solid()
pain_bar.fill.fore_color.rgb = LBLUE
pain_bar.line.color.rgb = LGRAY
pain_bar.line.width = Pt(0.5)

pain_texts = [
    "① AI推理成本爆炸：日均120万亿Token，16倍/年↑",
    "② 英伟达断供→供应链重构：H20退场，训练芯片缺口",
    "③ 视频生成算力不足：Seedance占80%算力，排队拥堵",
    "④ 国产芯片产能瓶颈：华为满产，寒武纪排至2027",
    "⑤ 存储芯片涨价：DRAM 250%+，HBM依赖海外",
]

pain_tb = add_textbox(slide2, MARGIN + Inches(0.1), pain_y + Inches(0.03),
                       Inches(12.5), Inches(0.6), "", font_size=Pt(9))
tf = pain_tb.text_frame
tf.word_wrap = True
for i, pt in enumerate(pain_texts):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = pt
    p.font.size = Pt(9)
    p.font.color.rgb = DGRAY
    p.font.name = FONT
    p.space_before = Pt(0)
    p.space_after = Pt(1)
    p.alignment = PP_ALIGN.LEFT


# ═══════════════════════════════════════════════════════════
#  第3页：AI高价值场景
# ═══════════════════════════════════════════════════════════

slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3)
add_title_bar(slide3, "AI高价值场景 — 三梯队聚焦昇腾可服务场景")

top_y3 = CONTENT_TOP
full_w = Inches(12.7)

def add_tier_table(slide, y, title, data, tier_color=BLUE):
    """添加一个梯队表格"""
    add_textbox(slide, MARGIN, y, full_w, Inches(0.3),
                title, font_size=Pt(11.5), bold=True, color=tier_color)
    rows = len(data) + 1
    cols = len(data[0])
    row_h = (Inches(1.3) if rows <= 5 else Inches(1.1)) / (rows - 1)
    tbl = slide.shapes.add_table(rows, cols, MARGIN, y + Inches(0.3),
                                  full_w, Inches(1.2)).table
    # 列宽
    col_ws = [Inches(2.8), Inches(3.3), Inches(3.3), Inches(3.3)]
    for i, w in enumerate(col_ws):
        tbl.columns[i].width = w

    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            set_cell_text(tbl.cell(r, c), val, font_size=Pt(9.5), bold=(r == 0))
    style_table_header(tbl)
    style_table_rows(tbl)
    return tbl

# P0 — 第一梯队
p0_data = [
    ["场景", "推理量级", "芯片需求", "机会点"],
    ["豆包大模型推理", "日均120万亿Token\n16倍/年", "NPU 25-35万颗\n+ASIC数百万颗", "昇腾950PR推理集群扩容\n高通ASIC部署"],
    ["Agent 7×24运行🆙", "单次请求→持续运行\n结构性质变", "NPU/ASIC常驻推理", "Agent运行时降本方案\n长记忆管理"],
    ["Seedance视频生成", "占市场80%\n排队拥堵", "视频推理专用NPU集群", "单卡推理效率极致优化\n降低单视频成本"],
    ["Trae AI编程推理", "Token半年+700%\n多模型聚合", "低延迟NPU推理", "<200ms TTFT\n推理链路优化"],
]
add_tier_table(slide3, top_y3, "第一梯队：推理算力密集型（P0 — 直接拉升NPU/ASIC需求）", p0_data)

# P1 — 第二梯队
p1_y = top_y3 + Inches(1.7)
p1_data = [
    ["场景", "算力特征", "芯片需求", "机会点"],
    ["推荐/广告LLM级训练", "千卡级MoE训练\n万级序列", "训练GPU ~7万颗", "昇腾910C训练集群\nCANN适配工具链"],
    ["火山引擎MaaS多模型", "49.5%调用量\n多模型混部", "推理NPU+GPU混合", "多模型路由调度\n训推一体平台"],
]
add_tier_table(slide3, p1_y, "第二梯队：训练+推理混合（P1）", p1_data)

# P2 — 第三梯队
p2_y = p1_y + Inches(1.5)
p2_data = [
    ["场景", "推理量级", "芯片需求", "机会点"],
    ["火山引擎汽车AI", "日均3000万+次\n700万+搭载", "端侧+云端混合NPU", "车载推理芯片方案"],
    ["飞书AI办公", "AaaS\n赛力斯174场景", "中轻量推理", "企业AI落地推理方案"],
]
add_tier_table(slide3, p2_y, "第三梯队：增量推理（P1-P2）", p2_data)


# ═══════════════════════════════════════════════════════════
#  第4页：算力需求与芯片供需
# ═══════════════════════════════════════════════════════════

slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4)
add_title_bar(slide4, "算力需求映射与推理芯片供需")

top_y4 = CONTENT_TOP
half_w = Inches(6.2)

# 左侧：业务→算力映射表
add_textbox(slide4, MARGIN, top_y4, half_w, Inches(0.3),
            "业务 → 算力映射", font_size=Pt(12), bold=True, color=BLUE)

map_data = [
    ["业务场景", "模型/负载", "类型", "当前规模"],
    ["豆包C端推理", "Seed-2.0-pro 200B MoE", "推理", "日均120万亿Token"],
    ["Agent生态", "ArkClaw/Coze 7×24", "推理", "单次→持续(结构变)"],
    ["Seedance视频", "双分支扩散Transformer", "推理", "市场80%算力"],
    ["Trae编程", "多模型聚合推理", "推理", "Token半年+700%"],
    ["推荐/广告训练", "OneTrans/HyFormer", "训练+推理", "万级序列LLM级"],
    ["火山引擎MaaS", "多模型混部", "推理", "全行业1944万亿Token"],
    ["汽车AI", "对话推理引擎", "推理", "日均3000万+次"],
    ["飞书AI", "aily Agent", "推理", "赛力斯174场景"],
]

tbl_map = slide4.shapes.add_table(len(map_data), 4, MARGIN, top_y4 + Inches(0.3),
                                   half_w, Inches(2.6)).table
tbl_map.columns[0].width = Inches(1.5)
tbl_map.columns[1].width = Inches(1.8)
tbl_map.columns[2].width = Inches(0.8)
tbl_map.columns[3].width = Inches(2.1)

for r, row_data in enumerate(map_data):
    for c, val in enumerate(row_data):
        set_cell_text(tbl_map.cell(r, c), val, font_size=Pt(8.5), bold=(r == 0))
style_table_header(tbl_map)
style_table_rows(tbl_map)

# 增长列
for r in range(1, len(map_data)):
    growth_map = {1: "↑↑↑", 2: "↑↑", 3: "↑↑↑", 4: "↑↑↑", 5: "↑↑", 6: "↑↑", 7: "↑", 8: "↑"}
    # add growth indicator via a small text
    growth_text = growth_map.get(r, "↑")
    set_cell_text(tbl_map.cell(r, 3),
                  tbl_map.cell(r, 3).text + f"   {growth_text}",
                  font_size=Pt(8.5))

# 右侧：推理芯片供需表
right_x4 = Inches(6.65)
add_textbox(slide4, right_x4, top_y4, half_w, Inches(0.3),
            "推理芯片供需", font_size=Pt(12), bold=True, color=BLUE)

chip_data = [
    ["类型", "供应商", "2026规模", "缺口", "动态"],
    ["推理NPU", "昇腾950PR", "25-35万颗", "产能预订满", "国产推理主力"],
    ["推理NPU", "寒武纪590", "30万颗", "排至2027", "Q1首盈>10亿"],
    ["推理ASIC", "高通", "数百万颗", "填补NPU缺口", "彭博5/26确认"],
    ["自研推理", "Ada", "10-35万颗", "起步", "三星6nm"],
    ["训练GPU", "昇腾910C", "~7万颗", "供应受限", "国产训练主力"],
    ["训练GPU", "AMD", "待确认", "替代英伟达", "新供应商"],
]

tbl_chip = slide4.shapes.add_table(len(chip_data), 5, right_x4, top_y4 + Inches(0.3),
                                    half_w, Inches(2.0)).table
tbl_chip.columns[0].width = Inches(1.0)
tbl_chip.columns[1].width = Inches(0.9)
tbl_chip.columns[2].width = Inches(1.1)
tbl_chip.columns[3].width = Inches(1.2)
tbl_chip.columns[4].width = Inches(2.0)

for r, row_data in enumerate(chip_data):
    for c, val in enumerate(row_data):
        set_cell_text(tbl_chip.cell(r, c), val, font_size=Pt(8.5), bold=(r == 0))
style_table_header(tbl_chip)
style_table_rows(tbl_chip)

# ── 下方：上下文长度暴增可视化 ──
ctx_y = top_y4 + Inches(3.15)
add_textbox(slide4, MARGIN, ctx_y, full_w, Inches(0.3),
            "上下文长度暴增 — 从ChatBot到Agent的算力隐形引擎",
            font_size=Pt(11.5), bold=True, color=BLUE)

ctx_card_y = ctx_y + Inches(0.3)
ctx_card_w = Inches(2.8)
ctx_card_h = Inches(1.1)
ctx_gap = Inches(0.35)
ctx_arrow_w = Inches(0.35)

ctx_data = [
    ("千Token级", "短对话\n单次问答即结束\nKV Cache小"),
    ("万Token级", "长文档\n多轮深度对话\nKV Cache线性增长"),
    ("百万Token级", "豆包v13.2.0\n超长上下文\nKV Cache平方级增长"),
    ("25h+持续运行", "Seed-2.0-lite\nAgent持续思考\n常驻满负荷推理"),
]

total_ctx_w = len(ctx_data) * ctx_card_w + (len(ctx_data) - 1) * (ctx_arrow_w + ctx_gap - Inches(0.05))
ctx_start_x = Inches(0.3) + (full_w - total_ctx_w) / 2

for i, (title, desc) in enumerate(ctx_data):
    cx = ctx_start_x + i * (ctx_card_w + ctx_arrow_w + ctx_gap - Inches(0.05))
    card = add_card(slide4, cx, ctx_card_y, ctx_card_w, ctx_card_h, fill_color=LBLUE, border_color=LGRAY)
    add_blue_left_border(slide4, cx, ctx_card_y, ctx_card_h, Inches(0.05))
    add_textbox(slide4, cx + Inches(0.12), ctx_card_y + Inches(0.08),
                ctx_card_w - Inches(0.15), Inches(0.3),
                title, font_size=Pt(10.5), bold=True, color=BLUE)
    add_textbox(slide4, cx + Inches(0.12), ctx_card_y + Inches(0.42),
                ctx_card_w - Inches(0.15), Inches(0.6),
                desc, font_size=Pt(8.5), color=DGRAY)

    # 箭头
    if i < len(ctx_data) - 1:
        ax = cx + ctx_card_w + Inches(0.02)
        ay = ctx_card_y + ctx_card_h/2 - Inches(0.08)
        arr = slide4.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, ax, ay, ctx_arrow_w, Inches(0.16))
        arr.fill.solid()
        arr.fill.fore_color.rgb = BLUE
        arr.line.fill.background()

# ── 底部总结文字 ──
summary_y = ctx_card_y + ctx_card_h + Inches(0.2)
add_textbox(slide4, MARGIN, summary_y, full_w, Inches(0.5),
            "💡 推理算力是最大、最确定的销售机会：日均120万亿Token + 16倍/年增速 + Agent 7×24持续运行 → "
            "推理芯片缺口千亿级市场。昇腾950PR(25-35万颗) + 寒武纪590(30万颗) + 高通ASIC(数百万颗) = 核心供应格局。",
            font_size=Pt(10), bold=True, color=DGRAY)


# ═══════════════════════════════════════════════════════════
#  第5页：决策链 + 洞察摘要
# ═══════════════════════════════════════════════════════════

slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5)
add_title_bar(slide5, "决策链与洞察摘要 — 触达路径与销售策略")

top_y5 = CONTENT_TOP

# ── 决策链横向卡片 ──
dc_data = [
    ("火山引擎总裁\n谭待", "芯片采购+AI Infra\n最终决策 ★★★★★"),
    ("火山引擎BU\n杨震原", "整体战略\n★★★★"),
    ("芯片团队\n施云峰/余红斌", "芯片选型+自研\n(含CPU) ★★★★"),
    ("Seed负责人\n吴永辉", "大模型技术方向\n★★★★"),
    ("CEO\n梁汝波", "终批\n★★★★★"),
]

dc_card_w = Inches(2.3)
dc_card_h = Inches(1.25)
dc_gap = Inches(0.2)
dc_arrow_w = Inches(0.25)

total_dc_w = len(dc_data) * dc_card_w + (len(dc_data) - 1) * dc_gap
dc_start_x = Inches(0.3) + (full_w - total_dc_w) / 2

for i, (name, desc) in enumerate(dc_data):
    dx = dc_start_x + i * (dc_card_w + dc_gap)
    card = add_card(slide5, dx, top_y5, dc_card_w, dc_card_h, fill_color=LBLUE, border_color=LGRAY)
    add_blue_left_border(slide5, dx, top_y5, dc_card_h, Inches(0.05))
    add_textbox(slide5, dx + Inches(0.1), top_y5 + Inches(0.08),
                dc_card_w - Inches(0.15), Inches(0.45),
                name, font_size=Pt(10.5), bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide5, dx + Inches(0.08), top_y5 + Inches(0.55),
                dc_card_w - Inches(0.16), Inches(0.65),
                desc, font_size=Pt(8.5), color=DGRAY, alignment=PP_ALIGN.CENTER)

    # 箭头
    if i < len(dc_data) - 1:
        ax = dx + dc_card_w + Inches(0.02)
        ay = top_y5 + dc_card_h/2 - Inches(0.07)
        arr = slide5.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, ax, ay, dc_arrow_w, Inches(0.14))
        arr.fill.solid()
        arr.fill.fore_color.rgb = BLUE
        arr.line.fill.background()

# 触达路径文字
path_y = top_y5 + dc_card_h + Inches(0.12)
add_textbox(slide5, MARGIN, path_y, full_w, Inches(0.3),
            "触达路径：火山引擎总裁团队 → Seed技术团队 → 芯片团队",
            font_size=Pt(10), bold=False, color=DGRAY, alignment=PP_ALIGN.CENTER)

# ── 洞察摘要卡片 ──
insight_y = path_y + Inches(0.4)
insight_card_w = full_w - Inches(0.0)
insight_card_h = Inches(1.1)
insight_gap = Inches(0.12)

insights = [
    ("① 推理算力是最大、最确定的销售机会——没有之一。",
     "字节从ChatBot到Agent的范式跃迁，叠加上下文窗口从千级Token扩展到百万级（豆包v13.2.0已支持百万Token窗口），使单次推理的KV Cache和计算量呈平方级增长。日均120万亿Token + 16倍/年增速 + Agent 7×24持续运行——三重杠杆叠加，推理芯片缺口是当前最紧迫需求。昇腾950PR（25-35万颗）+ 寒武纪590（30万颗）+ 高通ASIC（数百万颗）= 千亿级推理芯片市场。"),

    ("② 从ChatBot到Agent，上下文长度暴增是算力需求的隐形引擎。",
     "短对话（千Token级）→ 超长Agent任务（Seed-2.0-lite支持25小时+持续运行）→ 百万Token上下文窗口，每一次范式升级都带来推理算力的指数级增长。Agent不再是\"一次问答\"，而是\"持续思考+调用工具+长记忆\"——推理芯片需要从\"峰值弹性\"升级为\"常驻满负荷\"。"),

    ("③ AIGC-Agent是推理需求增长的新引擎。",
     "豆包3.45亿MAU + Seed-Thinking推理模型（200B MoE，思考链延长增加推理量）+ 视频生成占市场80% + Agent 7×24运行——这四场景的推理需求增长远超传统推荐系统，是未来2-3年推理芯片采购的主力驱动。"),

    ("④ 2026年850亿-4777亿AI芯片采购预算是中国最大的单客户AI芯片销售机会。",
     "谭待/杨震原（采购决策）→ 施云峰/余红斌（选型）→ 梁汝波（终批），决策链清晰、预算充裕、窗口期明确。"),
]

for i, (title, body) in enumerate(insights):
    iy = insight_y + i * (insight_card_h + insight_gap)
    card = add_card(slide5, MARGIN, iy, insight_card_w, insight_card_h,
                    fill_color=WHITE, border_color=LGRAY)
    add_blue_left_border(slide5, MARGIN, iy, insight_card_h, Inches(0.06))
    # 标题
    add_textbox(slide5, MARGIN + Inches(0.15), iy + Inches(0.06),
                insight_card_w - Inches(0.25), Inches(0.3),
                title, font_size=Pt(10.5), bold=True, color=BLUE)
    # 正文
    add_textbox(slide5, MARGIN + Inches(0.15), iy + Inches(0.35),
                insight_card_w - Inches(0.25), Inches(0.7),
                body, font_size=Pt(9), color=DGRAY)


# ═══════════════════════════════════════════════════════════
#  保存
# ═══════════════════════════════════════════════════════════

output_path = "/root/.openclaw/workspace/memory/ba-ta/字节/字节跳动_BA画像_Creater_v3.3.pptx"
prs.save(output_path)
print(f"✅ PPT已保存至: {output_path}")
print(f"   共 {len(prs.slides)} 页幻灯片")
