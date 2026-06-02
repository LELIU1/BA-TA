#!/usr/bin/env python3
"""
Generate 字节跳动 BA画像 PPT — 6 slides, 16:9, professional dark-header style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── Color Palette ──
DARK_HEADER  = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG     = RGBColor(0xF5, 0xF7, 0xFA)
BLUE_ACCENT  = RGBColor(0x3B, 0x82, 0xF6)
RED_ACCENT   = RGBColor(0xEF, 0x44, 0x44)
AMBER_ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
GREEN_ACCENT = RGBColor(0x10, 0xB9, 0x81)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT     = RGBColor(0x1F, 0x29, 0x37)
GRAY_TEXT     = RGBColor(0x6B, 0x72, 0x80)
CARD_BG      = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER  = RGBColor(0xE5, 0xE7, 0xEB)
TABLE_HDR_BG = RGBColor(0x1A, 0x1A, 0x2E)
TABLE_ALT_BG = RGBColor(0xF0, 0xF4, 0xFE)
CHIP_OVERSEAS = RGBColor(0x3B, 0x82, 0xF6)   # blue
CHIP_DOMESTIC = RGBColor(0xEF, 0x44, 0x44)   # red
CHIP_SELF     = RGBColor(0x10, 0xB9, 0x81)   # green

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helpers ──

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1), radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    if radius:
        # Adjust corner radius via XML
        sp = shape._element
        prstGeom = sp.find(qn('a:prstGeom'), sp.nsmap) or sp.find('.//' + qn('a:prstGeom'))
        if prstGeom is None:
            prstGeom = sp.find('.//' + qn('a:prstGeom'))
        # fallback: it should exist
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=Pt(12),
                color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name='PingFang SC',
                anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(2)
    # line spacing
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    tf.paragraphs[0].space_before = Pt(0)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, default_size=Pt(11),
                          default_color=DARK_TEXT, default_bold=False, font_name='PingFang SC',
                          alignment=PP_ALIGN.LEFT, line_spacing=1.3):
    """lines: list of (text, size, color, bold) or str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, size, color, bold = line, default_size, default_color, default_bold
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            color = line[2] if len(line) > 2 else default_color
            bold = line[3] if len(line) > 3 else default_bold
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(3)
    return txBox

def add_header_bar(slide, title_text, subtitle_text=None, page_num=None):
    """Dark header bar at top"""
    bar = add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.15), fill_color=DARK_HEADER)
    bar.line.fill.background()
    # Title
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.55),
                title_text, font_size=Pt(24), color=WHITE, bold=True)
    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.68), Inches(10), Inches(0.35),
                    subtitle_text, font_size=Pt(11), color=RGBColor(0x9C, 0xA3, 0xAF))
    # Page number
    if page_num:
        add_textbox(slide, Inches(11.8), Inches(0.35), Inches(1.2), Inches(0.45),
                    f'{page_num} / 6', font_size=Pt(11), color=RGBColor(0x9C, 0xA3, 0xAF),
                    alignment=PP_ALIGN.RIGHT)
    # Accent line under header
    add_rect(slide, Inches(0), Inches(1.15), prs.slide_width, Inches(0.04), fill_color=BLUE_ACCENT)

def add_footer(slide, source_text="数据来源：36氪 · Bloomberg · Reuters · 火山引擎官方 · IDC · QuestMobile · Sacra · 寒武纪财报 | 基于华为'看BA讲TA'方法论·察阶段 | 2026-06-01"):
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12), Inches(0.35),
                source_text, font_size=Pt(7), color=GRAY_TEXT)

def add_kpi_card(slide, left, top, width, height, value, label, color=BLUE_ACCENT):
    """A rounded card with big number and label"""
    card = add_rect(slide, left, top, width, height, fill_color=CARD_BG, border_color=CARD_BORDER, radius=True)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.55),
                value, font_size=Pt(18), color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.58), width - Inches(0.2), Inches(0.35),
                label, font_size=Pt(9), color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

def add_table(slide, left, top, col_widths, headers, rows, header_font=Pt(9), cell_font=Pt(8)):
    """Create a styled table. col_widths in Inches."""
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, Inches(0.35 * n_rows))
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = header_font
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = 'PingFang SC'
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = cell_font
                p.font.color.rgb = DARK_TEXT
                p.font.name = 'PingFang SC'
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape

def add_arrow_down(slide, left, top, width=Inches(0.04), height=Inches(0.3), color=BLUE_ACCENT):
    """Simple downward arrow shape"""
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow

def add_circle(slide, left, top, size, fill_color=BLUE_ACCENT):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill_color
    circle.line.fill.background()
    return circle


# ═══════════════════════════════════════════
# SLIDE 1: 客户基本画像 + 商业模式
# ═══════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1)
add_header_bar(slide1, '字节跳动 BA业务架构画像', '客户基本画像与商业模式 | 成熟期 + AI转型深水区', 1)

# ── Key KPI Cards row ──
kpi_data = [
    ('~1900亿$', '2025营收', BLUE_ACCENT),
    ('~500亿$', '净利润(AI投入致下滑>70%)', RED_ACCENT),
    ('~2000亿¥', '2026 AI Capex', AMBER_ACCENT),
    ('~4777亿¥', '内部上限($700亿)', RED_ACCENT),
    ('3.45亿', '豆包MAU', BLUE_ACCENT),
    ('120万亿/日', '日均Token消耗', BLUE_ACCENT),
]
kpi_left = Inches(0.4)
kpi_w = Inches(1.95)
kpi_h = Inches(0.95)
kpi_gap = Inches(0.1)
for i, (val, label, clr) in enumerate(kpi_data):
    add_kpi_card(slide1, kpi_left + i*(kpi_w + kpi_gap), Inches(1.4), kpi_w, kpi_h, val, label, clr)

# ── Core Business Matrix (left side) ──
matrix_left = Inches(0.4)
matrix_top = Inches(2.55)
matrix_w = Inches(7.0)

# Section title
add_textbox(slide1, matrix_left, matrix_top, Inches(4), Inches(0.3),
            '▎核心业务矩阵', font_size=Pt(14), color=DARK_HEADER, bold=True)

matrix_headers = ['类别', '核心产品', '规模/地位', 'AI算力关联']
matrix_cols = [Inches(1.1), Inches(2.1), Inches(1.9), Inches(1.9)]
matrix_rows = [
    ['内容与社区', '抖音 · TikTok · 头条 · 番茄小说', '抖音DAU 7亿+ TikTok DAU 10亿+', '推荐LLM化，推理↑↑'],
    ['创作工具', '剪映/CapCut · 即梦AI', '全球视频编辑TOP1, MAU 7994万', 'Seedance视频生成80%算力'],
    ['AI应用', '豆包 · Trae · 扣子Coze', '豆包MAU 3.45亿 Trae 600万', '日均120万亿Token推理'],
    ['企业服务', '火山引擎 · 飞书', 'MaaS 49.5%中国市场第一', 'Agent Plan订阅, AaaS底座'],
    ['电商', '抖音电商 · TikTok Shop', 'TikTok Shop GMV近千亿$', '海外AIDC核心驱动力'],
]
add_table(slide1, matrix_left, matrix_top + Inches(0.35), matrix_cols, matrix_headers, matrix_rows,
          header_font=Pt(9), cell_font=Pt(8))

# ── Right side: Revenue Structure + Regional Structure ──
right_x = Inches(7.7)
right_w = Inches(5.2)

# Revenue pie-style cards
add_textbox(slide1, right_x, Inches(2.55), Inches(4), Inches(0.3),
            '▎收入结构', font_size=Pt(14), color=DARK_HEADER, bold=True)

rev_data = [
    ('广告', '55-60%', BLUE_ACCENT, '抖音/TikTok信息流，点击率>行业30%'),
    ('电商', '15-20%', RED_ACCENT, '抖音电商+TikTok Shop GMV近千亿$'),
    ('云+AI', '5-10%', GREEN_ACCENT, '火山引擎MaaS 49.5%第一，豆包付费'),
    ('游戏/教育/其它', '5-10%', GRAY_TEXT, '朝夕光年·大力教育'),
]
for j, (name, pct, clr, desc) in enumerate(rev_data):
    y = Inches(2.95) + j * Inches(0.75)
    c = add_rect(slide1, right_x, y, right_w, Inches(0.65), fill_color=CARD_BG, border_color=CARD_BORDER, radius=True)
    # colored left bar
    add_rect(slide1, right_x + Inches(0.03), y + Inches(0.08), Inches(0.06), Inches(0.49), fill_color=clr)
    add_textbox(slide1, right_x + Inches(0.2), y + Inches(0.05), Inches(1.2), Inches(0.3),
                pct, font_size=Pt(18), color=clr, bold=True)
    add_textbox(slide1, right_x + Inches(1.35), y + Inches(0.1), Inches(1.5), Inches(0.25),
                name, font_size=Pt(11), color=DARK_TEXT, bold=True)
    add_textbox(slide1, right_x + Inches(0.2), y + Inches(0.33), Inches(4.5), Inches(0.25),
                desc, font_size=Pt(8), color=GRAY_TEXT)

# ── Regional + Value Proposition (bottom row) ──
add_textbox(slide1, right_x, Inches(5.95), Inches(4), Inches(0.25),
            '▎地域与价值主张', font_size=Pt(14), color=DARK_HEADER, bold=True)

# Regional cards
regions = [
    ('国内', '~70%', '增速~20%', BLUE_ACCENT),
    ('海外', '~30%(+5pp)', '增速~50%', RED_ACCENT),
]
for k, (name, share, growth, clr) in enumerate(regions):
    rx = right_x + k * Inches(2.6)
    add_rect(slide1, rx, Inches(6.3), Inches(2.4), Inches(0.55), fill_color=CARD_BG, border_color=CARD_BORDER, radius=True)
    add_textbox(slide1, rx + Inches(0.12), Inches(6.33), Inches(1.0), Inches(0.22),
                name, font_size=Pt(10), color=DARK_TEXT, bold=True)
    add_textbox(slide1, rx + Inches(0.12), Inches(6.55), Inches(0.8), Inches(0.22),
                share, font_size=Pt(14), color=clr, bold=True)
    add_textbox(slide1, rx + Inches(0.9), Inches(6.55), Inches(1.2), Inches(0.22),
                growth, font_size=Pt(9), color=GRAY_TEXT)

add_footer(slide1)


# ═══════════════════════════════════════════
# SLIDE 2: 关键AI基建供应商
# ═══════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)
add_header_bar(slide2, '关键AI基建供应商', '海外芯片供应链 · 国产芯片供应链 · 自研芯片 | 从"英伟达独供"到"高通+AMD+Broadcom"新三角', 2)

# Three sections side by side
section_data = [
    {
        'title': '🌍 海外芯片供应链',
        'color': CHIP_OVERSEAS,
        'items': [
            ('🔴 高通', 'AI ASIC + 联合制造', '数百万颗', '推理ASIC主力'),
            ('AMD', 'GPU（MI系列）', '替代英伟达', '训练+推理GPU'),
            ('Broadcom', '网络芯片/ASIC IP', 'AIDC网络+IP合作', '网络基础设施'),
        ],
        'note': '英伟达份额→~8%，实质退场'
    },
    {
        'title': '🇨🇳 国产芯片供应链',
        'color': CHIP_DOMESTIC,
        'items': [
            ('华为昇腾', '910C+950PR', '25-35万颗950PR', '训练+推理主力'),
            ('寒武纪', '思元590', '30万颗,排至2027', 'Day-0适配DeepSeek-V4'),
        ],
        'note': '华为~300-400亿 | 寒武纪Q1首盈>10亿'
    },
    {
        'title': '🔧 自研芯片',
        'color': CHIP_SELF,
        'items': [
            ('Ada/SeedChip', '三星6nm推理', '目标10-35万颗', '2026量产推进'),
            ('自研CPU 🆕', 'ARM+RISC-V', '早期设计阶段', '路透5/28曝光'),
            ('高通联合', 'ASIC设计+制造', '数百万颗', '订单推进中'),
        ],
        'note': 'CPU双路线并行，供应链多元化'
    },
]

sec_w = Inches(3.85)
sec_gap = Inches(0.35)
start_x = Inches(0.4)
start_y = Inches(1.45)

for i, sec in enumerate(section_data):
    sx = start_x + i * (sec_w + sec_gap)
    # Section card
    card = add_rect(slide2, sx, start_y, sec_w, Inches(5.35), fill_color=CARD_BG, border_color=CARD_BORDER, radius=True)
    # Colored top bar
    add_rect(slide2, sx, start_y, sec_w, Inches(0.5), fill_color=sec['color'], radius=False)
    # The rounded card doesn't play well with top-only radius; overlay a text box
    add_textbox(slide2, sx + Inches(0.15), start_y + Inches(0.05), sec_w - Inches(0.3), Inches(0.4),
                sec['title'], font_size=Pt(14), color=WHITE, bold=True)

    # Supplier cards
    for j, (name, product, scale, pos) in enumerate(sec['items']):
        iy = start_y + Inches(0.7) + j * Inches(1.35)
        item_card = add_rect(slide2, sx + Inches(0.12), iy, sec_w - Inches(0.24), Inches(1.2),
                             fill_color=LIGHT_BG, border_color=CARD_BORDER, radius=True)
        # Name
        add_textbox(slide2, sx + Inches(0.25), iy + Inches(0.06), sec_w - Inches(0.5), Inches(0.3),
                    name, font_size=Pt(12), color=DARK_TEXT, bold=True)
        # Product
        add_textbox(slide2, sx + Inches(0.25), iy + Inches(0.38), sec_w - Inches(0.5), Inches(0.22),
                    f'提供: {product}', font_size=Pt(8.5), color=BLUE_ACCENT)
        # Scale
        add_textbox(slide2, sx + Inches(0.25), iy + Inches(0.6), sec_w - Inches(0.5), Inches(0.22),
                    f'规模: {scale}', font_size=Pt(9), color=DARK_TEXT, bold=True)
        # Position
        add_textbox(slide2, sx + Inches(0.25), iy + Inches(0.85), sec_w - Inches(0.5), Inches(0.22),
                    f'定位: {pos}', font_size=Pt(8), color=GRAY_TEXT)

    # Note at bottom of section
    add_textbox(slide2, sx + Inches(0.15), start_y + Inches(4.85), sec_w - Inches(0.3), Inches(0.35),
                f'💡 {sec["note"]}', font_size=Pt(8.5), color=AMBER_ACCENT, bold=True)

add_footer(slide2)


# ═══════════════════════════════════════════
# SLIDE 3: 核心业务流程 — 双轨制
# ═══════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)
add_header_bar(slide3, '核心业务流程：双轨制AI价值链路', '左轨：搜索推荐AI闭环（传统核心）  |  右轨：AIGC-Agent应用闭环（新增长极⭐主线）', 3)

# ── Left Track ──
lt_left = Inches(0.35)
lt_top = Inches(1.45)
lt_w = Inches(5.8)

# Left track title
left_track_bg = add_rect(slide3, lt_left, lt_top, lt_w, Inches(5.4), fill_color=CARD_BG, border_color=BLUE_ACCENT, border_width=Pt(2), radius=True)
add_rect(slide3, lt_left, lt_top, lt_w, Inches(0.5), fill_color=BLUE_ACCENT)
add_textbox(slide3, lt_left + Inches(0.15), lt_top + Inches(0.06), lt_w - Inches(0.3), Inches(0.38),
            ' 左轨：搜索推荐AI闭环', font_size=Pt(14), color=WHITE, bold=True)

# Flow boxes
left_flow = [
    ('抖音/TikTok/头条\n用户行为', '海量行为数据'),
    ('推荐算法 LLM化\nRankMixer→HyFormer', '万级序列长度'),
    ('广告投放+电商转化\n核心收入 ~70%', '现金流引擎'),
    ('持续投入\nAI基础设施', 'Capex >2000亿元'),
    ('更强大推荐模型\n更长序列·更大参数', '体验升级飞轮'),
]
box_h = Inches(0.72)
box_gap_v = Inches(0.18)
flow_start_y = lt_top + Inches(0.65)
for j, (title, sub) in enumerate(left_flow):
    by = flow_start_y + j * (box_h + box_gap_v + Inches(0.1))
    bx = lt_left + Inches(0.25)
    bw = lt_w - Inches(0.5)
    add_rect(slide3, bx, by, bw, box_h, fill_color=LIGHT_BG, border_color=CARD_BORDER, radius=True)
    add_textbox(slide3, bx + Inches(0.15), by + Inches(0.08), bw - Inches(0.3), Inches(0.38),
                title, font_size=Pt(10), color=DARK_TEXT, bold=True)
    add_textbox(slide3, bx + Inches(0.15), by + Inches(0.46), bw - Inches(0.3), Inches(0.22),
                sub, font_size=Pt(8), color=GRAY_TEXT)
    # Arrow between boxes (except last)
    if j < len(left_flow) - 1:
        ax = bx + bw/2 - Inches(0.06)
        add_arrow_down(slide3, ax, by + box_h + Inches(0.02), height=Inches(0.15))

# Left bottom: "封闭飞轮" indicator
add_rect(slide3, lt_left + Inches(0.5), flow_start_y + 5*(box_h + box_gap_v + 0.1) + Inches(0.05),
         lt_w - Inches(1), Inches(0.4), fill_color=BLUE_ACCENT)
add_textbox(slide3, lt_left + Inches(0.5), flow_start_y + 5*(box_h + box_gap_v + 0.1) + Inches(0.1),
         lt_w - Inches(1), Inches(0.3), '↻ 封闭增长飞轮', font_size=Pt(11), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Feature note
add_textbox(slide3, lt_left + Inches(0.15), lt_top + Inches(4.95), lt_w - Inches(0.3), Inches(0.35),
            '⚙ 芯片需求: 训练GPU（昇腾910C+AMD）为主，模型参数向LLM级演进',
            font_size=Pt(8.5), color=BLUE_ACCENT, bold=True)

# ── Right Track ──
rt_left = Inches(6.65)
rt_top = Inches(1.45)
rt_w = Inches(6.3)

right_track_bg = add_rect(slide3, rt_left, rt_top, rt_w, Inches(5.4), fill_color=CARD_BG, border_color=RED_ACCENT, border_width=Pt(2), radius=True)
add_rect(slide3, rt_left, rt_top, rt_w, Inches(0.5), fill_color=RED_ACCENT)
add_textbox(slide3, rt_left + Inches(0.15), rt_top + Inches(0.06), rt_w - Inches(0.3), Inches(0.38),
            ' 右轨：AIGC-Agent应用闭环 ⭐ 主线', font_size=Pt(14), color=WHITE, bold=True)

# Top node: Seed model
add_rect(slide3, rt_left + Inches(1.5), rt_top + Inches(0.65), Inches(3.3), Inches(0.55),
         fill_color=RED_ACCENT)
add_textbox(slide3, rt_left + Inches(1.5), rt_top + Inches(0.68), Inches(3.3), Inches(0.5),
            '豆包大模型 Seed-2.0-pro/lite 全模态', font_size=Pt(11), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_arrow_down(slide3, rt_left + Inches(3.0), rt_top + Inches(1.2), height=Inches(0.15), color=RED_ACCENT)

# App matrix - 5 cards in 2 rows
app_items = [
    ('🤖 豆包App', 'MAU 3.45亿\nv13.2 百万Token窗口\n买前问豆包电商闭环'),
    ('🎬 即梦+剪映\nSeedance 2.0', 'AI短剧日更>1100部\n视频生成算力80%\n火山剧创1.0'),
    ('💻 Trae AI编程', '600万开发者\nSOLO模式全流程\nToken半年+700%'),
    ('🧩 火山引擎MaaS', '调用量份额49.5%\nAgent Plan多模型\n汽车AI 700万+搭载'),
    ('🏭 企业Agent运营', 'Coze/ArkClaw 7×24\nHiAgent降本63%\n飞书aily'),
]
app_start_y = rt_top + Inches(1.5)
for k, (name, desc) in enumerate(app_items):
    col = k % 3
    row = k // 3
    ax = rt_left + Inches(0.2) + col * Inches(2.02)
    ay = app_start_y + row * Inches(1.55)
    add_rect(slide3, ax, ay, Inches(1.88), Inches(1.38), fill_color=LIGHT_BG, border_color=CARD_BORDER, radius=True)
    add_textbox(slide3, ax + Inches(0.08), ay + Inches(0.05), Inches(1.72), Inches(0.4),
                name, font_size=Pt(9.5), color=DARK_TEXT, bold=True)
    add_textbox(slide3, ax + Inches(0.08), ay + Inches(0.48), Inches(1.72), Inches(0.85),
                desc, font_size=Pt(7.5), color=GRAY_TEXT)

# Bottom: Token consumption
add_rect(slide3, rt_left + Inches(0.4), rt_top + Inches(4.65), Inches(5.5), Inches(0.55),
         fill_color=RED_ACCENT)
add_textbox(slide3, rt_left + Inches(0.4), rt_top + Inches(4.7), Inches(5.5), Inches(0.5),
            '日均120万亿Token推理 + Agent 7×24持续运行 → 驱动AI基础设施扩容',
            font_size=Pt(11), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Feature note
add_textbox(slide3, rt_left + Inches(0.15), rt_top + Inches(5.3), rt_w - Inches(0.3), Inches(0.35),
            '⚙ 芯片需求: 推理NPU（昇腾950PR+寒武纪590）+ 定制ASIC（高通）+ 自研（Ada/CPU）',
            font_size=Pt(8.5), color=RED_ACCENT, bold=True)

# ── Bottom: Pain Points ──
pain_y = Inches(6.9)
add_textbox(slide3, Inches(0.4), pain_y, Inches(12.5), Inches(0.3),
            '▎业务痛点（聚焦芯片销售视角）',
            font_size=Pt(11), color=DARK_HEADER, bold=True)
pain_items = [
    ('🔴 推理成本爆炸', '日均120万亿Token，年增16倍'),
    ('🔴 英伟达断供→重构', 'H20停供，存量→~8%'),
    ('🟡 视频算力不足', 'Seedance 80%算力，排队'),
    ('🟡 国产芯片产能瓶颈', '华为满/寒武纪排至2027'),
    ('🟡 存储芯片涨价', 'DRAM 250%+，HBM依赖海外'),
]
for m, (title, desc) in enumerate(pain_items):
    px = Inches(0.4) + m * Inches(2.5)
    add_rect(slide3, px, pain_y + Inches(0.28), Inches(2.3), Inches(0.35),
             fill_color=CARD_BG, border_color=CARD_BORDER, radius=True)
    add_textbox(slide3, px + Inches(0.08), pain_y + Inches(0.29), Inches(2.14), Inches(0.32),
                f'{title}: {desc}', font_size=Pt(7), color=DARK_TEXT)

# Override footer (positioned manually for this slide)
add_textbox(slide3, Inches(0.5), Inches(0.0), Inches(0.1), Inches(0.1), '', font_size=Pt(1))


# ═══════════════════════════════════════════
# SLIDE 4: AI高价值场景
# ═══════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)
add_header_bar(slide4, 'AI高价值场景：三梯队矩阵', '聚焦昇腾可服务场景 | P0推理密集型 → P1训推混合 → P2增量推理', 4)

# ── P0 Table ──
p0_y = Inches(1.45)
add_rect(slide4, Inches(0.35), p0_y, Inches(1.3), Inches(0.4), fill_color=RED_ACCENT)
add_textbox(slide4, Inches(0.35), p0_y + Inches(0.03), Inches(1.3), Inches(0.35),
            'P0 推理密集', font_size=Pt(11), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide4, Inches(1.7), p0_y + Inches(0.05), Inches(7), Inches(0.3),
            '直接拉升NPU/ASIC需求 — 最大销售机会', font_size=Pt(10), color=RED_ACCENT, bold=True)

p0_headers = ['场景', '推理量级', '芯片需求', '机会点']
p0_cols = [Inches(2.3), Inches(2.8), Inches(2.8), Inches(4.5)]
p0_rows = [
    ['豆包大模型推理', '日均120万亿Token\n16倍/年增长', 'NPU 25-35万颗\n+ ASIC数百万颗', '昇腾950PR推理集群扩容\n高通ASIC部署'],
    ['Agent 7×24运行 🆙', '单次请求→持续运行\n结构性质变', 'NPU/ASIC 常驻推理', 'Agent运行时降本方案\n长记忆管理优化'],
    ['Seedance视频生成', '占市场80%算力\n排队拥堵', '视频推理专用NPU集群', '单卡推理效率极致优化\n降低单视频成本'],
    ['Trae AI编程推理', 'Token半年+700%\n多模型聚合', '低延迟NPU推理', '<200ms TTFT推理链路优化'],
]
add_table(slide4, Inches(0.35), p0_y + Inches(0.5), p0_cols, p0_headers, p0_rows,
          header_font=Pt(9), cell_font=Pt(8))

# ── P1 Table ──
p1_y = Inches(3.65)
add_rect(slide4, Inches(0.35), p1_y, Inches(1.3), Inches(0.4), fill_color=AMBER_ACCENT)
add_textbox(slide4, Inches(0.35), p1_y + Inches(0.03), Inches(1.3), Inches(0.35),
            'P1 训推混合', font_size=Pt(11), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

p1_cols = [Inches(2.8), Inches(2.8), Inches(2.8), Inches(4.0)]
p1_rows = [
    ['推荐/广告LLM级训练', '千卡级MoE训练\n万级序列长度', '训练GPU ~7万颗', '昇腾910C训练集群\nCANN适配工具链'],
    ['火山引擎MaaS多模型', '49.5%调用量\n多模型混部', '推理NPU + GPU混合', '多模型路由调度\n训推一体平台'],
]
add_table(slide4, Inches(0.35), p1_y + Inches(0.5), p1_cols, p0_headers, p1_rows,
          header_font=Pt(9), cell_font=Pt(8))

# ── P2 Table ──
p2_y = Inches(4.95)
add_rect(slide4, Inches(0.35), p2_y, Inches(1.3), Inches(0.4), fill_color=BLUE_ACCENT)
add_textbox(slide4, Inches(0.35), p2_y + Inches(0.03), Inches(1.3), Inches(0.35),
            'P1-P2 增量', font_size=Pt(11), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

p2_cols = [Inches(2.8), Inches(2.8), Inches(2.8), Inches(4.0)]
p2_rows = [
    ['火山引擎汽车AI', '日均3000万+次座舱交互\n700万+搭载', '端侧+云端混合NPU', '车载推理芯片方案'],
    ['飞书AI办公', 'AaaS，赛力斯174场景', '中轻量推理', '企业AI落地推理方案'],
]
add_table(slide4, Inches(0.35), p2_y + Inches(0.5), p2_cols, p0_headers, p2_rows,
          header_font=Pt(9), cell_font=Pt(8))

add_footer(slide4)


# ═══════════════════════════════════════════
# SLIDE 5: 算力需求与芯片供应 + 上下文长度暴增可视化
# ═══════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)
add_header_bar(slide5, '场景算力需求与芯片供应', '上下文长度暴增：千Token → 万Token → 百万Token | 推理芯片缺口是当前最紧迫需求', 5)

# ── Left: 算力需求映射表 ──
add_textbox(slide5, Inches(0.4), Inches(1.45), Inches(6), Inches(0.3),
            '▎业务场景算力需求映射', font_size=Pt(13), color=DARK_HEADER, bold=True)

demand_headers = ['业务场景', '模型/负载', '类型', '当前规模', '增长']
demand_cols = [Inches(2.0), Inches(2.6), Inches(0.9), Inches(2.3), Inches(0.9)]
demand_rows = [
    ['豆包C端推理', 'Seed-2.0-pro(200B MoE)\n+Thinking-v1.5', '推理', '日均120万亿Token', '↑↑↑'],
    ['Agent生态', 'ArkClaw/Coze 7×24', '推理', '结构性质变', '↑↑'],
    ['Seedance视频', '双分支扩散Transformer', '推理', '市场80%，排队', '↑↑↑'],
    ['Trae编程', '多模型聚合推理', '推理', 'Token半年+700%', '↑↑↑'],
    ['推荐/广告训练', 'OneTrans/HyFormer', '训练+推理', '万级序列LLM级', '↑↑'],
    ['火山引擎MaaS', '多模型混部', '推理', '1944万亿Token全行业', '↑↑'],
    ['汽车AI', '对话推理引擎', '推理', '日均3000万+次', '↑'],
    ['飞书AI', 'aily Agent', '推理', '赛力斯174场景', '↑'],
]
add_table(slide5, Inches(0.35), Inches(1.85), demand_cols, demand_headers, demand_rows,
          header_font=Pt(8), cell_font=Pt(7.5))

# ── Right: 芯片供需表 ──
add_textbox(slide5, Inches(9.0), Inches(1.45), Inches(4.2), Inches(0.3),
            '▎芯片供需（推理芯片聚焦）', font_size=Pt(13), color=DARK_HEADER, bold=True)

supply_headers = ['类型', '供应商', '规模', '缺口']
supply_cols = [Inches(1.0), Inches(1.1), Inches(1.2), Inches(0.9)]
supply_rows = [
    ['推理NPU', '昇腾950PR', '25-35万颗', '产能满'],
    ['推理NPU', '寒武纪590', '30万颗', '排至2027'],
    ['推理ASIC', '高通', '数百万颗', '填补缺口'],
    ['自研推理', 'Ada', '10-35万颗', '起步'],
    ['训练GPU', '昇腾910C', '~7万颗', '供应受限'],
    ['训练GPU', 'AMD', '待确认', '替代英伟达'],
]
add_table(slide5, Inches(9.0), Inches(1.85), supply_cols, supply_headers, supply_rows,
          header_font=Pt(8), cell_font=Pt(7.5))

# ── Context Length Visualization ──
ctx_y = Inches(4.9)
add_rect(slide5, Inches(0.35), ctx_y, Inches(12.6), Inches(2.0), fill_color=CARD_BG, border_color=BLUE_ACCENT, border_width=Pt(1.5), radius=True)
add_textbox(slide5, Inches(0.6), ctx_y + Inches(0.08), Inches(6), Inches(0.3),
            '📐 上下文长度暴增：推理算力需求指数级增长的隐形引擎',
            font_size=Pt(13), color=DARK_HEADER, bold=True)

# Three context level cards
ctx_levels = [
    ('千Token级', '传统搜索推荐\n短对话理解', '基线', BLUE_ACCENT, Inches(3.5)),
    ('万Token级', '长文档分析\n多轮Agent任务', '10x增长', AMBER_ACCENT, Inches(7.5)),
    ('百万Token级', '豆包v13.2.0窗口\n25h+持续Agent运行', '1000x增长', RED_ACCENT, Inches(11.5)),
]
for i, (label, desc, growth, clr, arrow_x) in enumerate(ctx_levels):
    cx = Inches(0.7) + i * Inches(4.1)
    cw = Inches(2.5)
    # Card
    add_rect(slide5, cx, ctx_y + Inches(0.5), cw, Inches(1.2), fill_color=LIGHT_BG, border_color=clr, border_width=Pt(1.5), radius=True)
    # Colored top
    add_rect(slide5, cx, ctx_y + Inches(0.5), cw, Inches(0.4), fill_color=clr)
    add_textbox(slide5, cx + Inches(0.08), ctx_y + Inches(0.54), cw - Inches(0.16), Inches(0.32),
                label, font_size=Pt(13), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide5, cx + Inches(0.08), ctx_y + Inches(0.98), cw - Inches(0.16), Inches(0.28),
                desc, font_size=Pt(8.5), color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide5, cx + Inches(0.08), ctx_y + Inches(1.28), cw - Inches(0.16), Inches(0.28),
                f'算力指数: {growth}', font_size=Pt(9), color=clr, bold=True, alignment=PP_ALIGN.CENTER)

    # Arrow between levels
    if i < 2:
        arrow = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, cx + cw + Inches(0.1),
                                         ctx_y + Inches(0.98), Inches(1.3), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY_TEXT
        arrow.line.fill.background()
        add_textbox(slide5, cx + cw + Inches(0.1), ctx_y + Inches(1.24), Inches(1.3), Inches(0.2),
                    f'→ {growth}', font_size=Pt(7), color=RED_ACCENT if i == 1 else AMBER_ACCENT,
                    bold=True, alignment=PP_ALIGN.CENTER)

add_footer(slide5)


# ═══════════════════════════════════════════
# SLIDE 6: 决策链 + 洞察摘要
# ═══════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)
add_header_bar(slide6, '客户决策链与洞察摘要', '决策链清晰·预算充裕·窗口期明确 | 2026年最大单客户AI芯片销售机会', 6)

# ── Left: Decision Chain ──
add_textbox(slide6, Inches(0.4), Inches(1.45), Inches(4), Inches(0.3),
            '▎客户决策链与触达路径', font_size=Pt(13), color=DARK_HEADER, bold=True)

decision_members = [
    ('谭待', '火山引擎总裁', '★★★★★', '芯片采购+AI Infra\n最终决策', BLUE_ACCENT),
    ('梁汝波', 'CEO', '★★★★★', '终批', RED_ACCENT),
    ('施云峰/余红斌', '芯片团队', '★★★★', '芯片选型+自研\n含CPU路线', AMBER_ACCENT),
    ('吴永辉', 'Seed负责人', '★★★★', '大模型技术方向', GREEN_ACCENT),
    ('杨震原', '火山引擎BU负责人', '★★★★', '整体AI战略', BLUE_ACCENT),
]
for i, (name, role, influence, focus, clr) in enumerate(decision_members):
    dy = Inches(1.9) + i * Inches(0.82)
    # Colored left bar
    add_rect(slide6, Inches(0.55), dy, Inches(0.06), Inches(0.65), fill_color=clr)
    # Card
    add_rect(slide6, Inches(0.7), dy, Inches(5.3), Inches(0.65), fill_color=CARD_BG, border_color=CARD_BORDER, radius=True)
    # Star and name
    add_textbox(slide6, Inches(0.85), dy + Inches(0.05), Inches(1.8), Inches(0.28),
                f'{name}  {influence}', font_size=Pt(11), color=DARK_TEXT, bold=True)
    add_textbox(slide6, Inches(0.85), dy + Inches(0.33), Inches(2.5), Inches(0.25),
                role, font_size=Pt(8), color=GRAY_TEXT)
    add_textbox(slide6, Inches(3.6), dy + Inches(0.1), Inches(2.3), Inches(0.45),
                focus, font_size=Pt(8.5), color=DARK_TEXT)

# Touch path
add_textbox(slide6, Inches(0.7), Inches(6.05), Inches(5.5), Inches(0.3),
            '触达路径：火山引擎总裁团队 → Seed技术团队 → 芯片团队',
            font_size=Pt(9.5), color=BLUE_ACCENT, bold=True)

# ── Right: Insight Cards ──
add_textbox(slide6, Inches(6.8), Inches(1.45), Inches(6), Inches(0.3),
            '▎洞察摘要：4条核心洞察', font_size=Pt(13), color=DARK_HEADER, bold=True)

insights = [
    ('① 推理算力是最大最确定的销售机会',
     '从ChatBot到Agent的范式跃迁 + 上下文窗口百万Token → KV Cache平方级增长。日均120万亿Token + 16倍/年 + Agent 7×24 = 千亿级推理芯片市场。',
     RED_ACCENT, '🔥'),
    ('② 上下文暴增是算力需求的隐形引擎',
     '千Token→万Token→百万Token，每次范式升级带来推理算力指数级增长。Agent不再是一次问答，而是持续思考+调用工具+长记忆 = 芯片从峰值弹性升级为常驻满负荷。',
     AMBER_ACCENT, '⚡'),
    ('③ AIGC-Agent是推理需求增长新引擎',
     '豆包3.45亿MAU + Seed-Thinking(200B MoE) + 视频生成80%算力 + Agent 7×24运行 — 四场景需求远超传统推荐系统，未来2-3年推理芯片采购主力驱动。',
     BLUE_ACCENT, '🚀'),
    ('④ 2026年850亿-4777亿$采购预算：中国最大AI芯片销售机会',
     '谭待（采购决策）→ 施云峰/余红斌（选型）→ 梁汝波（终批）— 决策链清晰、预算充裕、窗口期明确。',
     GREEN_ACCENT, '💎'),
]
for i, (title, desc, clr, icon) in enumerate(insights):
    iy = Inches(1.88) + i * Inches(1.28)
    # Card
    add_rect(slide6, Inches(6.8), iy, Inches(6.15), Inches(1.12), fill_color=CARD_BG, border_color=clr, border_width=Pt(1.5), radius=True)
    # Colored accent dot
    add_circle(slide6, Inches(6.95), iy + Inches(0.12), Inches(0.18), fill_color=clr)
    # Title
    add_textbox(slide6, Inches(7.2), iy + Inches(0.06), Inches(5.5), Inches(0.3),
                f'{icon}  {title}', font_size=Pt(10.5), color=clr, bold=True)
    # Description
    add_textbox(slide6, Inches(6.95), iy + Inches(0.42), Inches(5.8), Inches(0.58),
                desc, font_size=Pt(7.5), color=DARK_TEXT)

add_footer(slide6)


# ── Save ──
output_path = '/root/.openclaw/workspace/memory/ba-ta/字节/字节跳动_BA画像_Creater.pptx'
prs.save(output_path)
print(f'✅ PPT saved to: {output_path}')

import os
size_mb = os.path.getsize(output_path) / 1024.0
print(f'📦 File size: {size_mb:.1f} KB')
